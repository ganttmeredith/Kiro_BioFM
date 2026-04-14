"""Generate PowerPoint deck for First Non-responders (AWS + Kiro narrative)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
LOGO = ROOT / "first_nonresponders_logo.png"
OUT = ROOT / "First_Nonresponders_AWS_Kiro_Deck.pptx"


def add_title_slide(prs, title: str, subtitle: str):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.55), Inches(0.35), width=Inches(4.2))
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(2.15), Inches(8.9), Inches(1.2))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p2 = tx.text_frame.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(18)
    p2.space_before = Pt(12)


def add_bullets(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)


def add_two_column(prs, title: str, left_title: str, left_bullets: list[str], right_title: str, right_bullets: list[str]):
    layout = prs.slide_layouts[5]  # title only
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(4.5), Inches(5.5))
    lf = left.text_frame
    h = lf.paragraphs[0]
    h.text = left_title
    h.font.bold = True
    h.font.size = Pt(22)
    for b in left_bullets:
        p = lf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)
    right = slide.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(4.5), Inches(5.5))
    rf = right.text_frame
    h2 = rf.paragraphs[0]
    h2.text = right_title
    h2.font.bold = True
    h2.font.size = Pt(22)
    for b in right_bullets:
        p = rf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(16)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "First Non-responders",
        "A clinician-led walkthrough: spot non-responders in the dashboard, then ask what matters",
    )

    add_bullets(
        prs,
        "Scenario: Dr. Alvarez in clinic prep",
        [
            "Head-and-neck cohort with mixed response—some patients progress earlier than expected.",
            "She wants two things: who looks like a non-responder in this filtered set, and why.",
            "The dashboard is her single workspace: explore morphology, match similar patients, then interrogate features in chat.",
            "No black box: every step ties back to cohort filters, slides, and explainable views she can share at tumor board.",
        ],
    )

    add_bullets(
        prs,
        "Walkthrough — our Cloudscape dashboard",
        [
            "Data — confirm the cohort loaded; review summary and metadata before interpreting plots.",
            "Explore — UMAP / embedding view with filters; visually isolate outliers and suspected non-responders.",
            "Morphology Groups — cluster patterns across slides to see who behaves like whom at tissue level.",
            "Patient Matcher — pick a concerning case; retrieve morphologically similar patients for comparison.",
            "Chat — natural-language questions: “What separates these patients?” with tool-backed answers and chart artifacts.",
        ],
    )

    add_bullets(
        prs,
        "Bedrock AgentCore (why judges should care)",
        [
            "AgentCore is the packaged runtime for our agent: BedrockAgentCoreApp + a single handler entrypoint.",
            "Same contract the UI expects: prompt or message history plus app_context (filters, patient focus, retrieval weights).",
            "Operational story: deploy once, invoke reliably—ideal for demos and production-style agent hosting on AWS.",
            "Under the hood, the handler streams the same Strands + Bedrock model path clinicians see in Chat.",
            "Bottom line: the dashboard is the face; AgentCore is the scalable back end for the assistant.",
        ],
    )

    add_two_column(
        prs,
        "AWS architecture (AgentCore first)",
        "Bedrock AgentCore",
        [
            "Handler returns structured results: assistant text plus optional artifacts from tool calls.",
            "Keeps prompts, tools, and model wiring in one deployable unit—not a one-off script.",
            "Pairs naturally with Bedrock: managed models, credentials, and regional deployment.",
        ],
        "Bedrock & platform",
        [
            "Foundation models on Amazon Bedrock (e.g. Claude) for reasoning and streaming replies.",
            "BEDROCK_MODEL_ID and AWS_DEFAULT_REGION configure model and region consistently.",
            "boto3 bedrock-runtime for health checks and report-style generation where needed.",
            "Large slides / assets: S3-style paths when datasets exceed local or container limits.",
        ],
    )

    add_bullets(
        prs,
        "How we used Kiro",
        [
            "Shipped faster: agent prompts, tool schemas, and AgentCore handler evolved in one AI-native workspace.",
            "Steering + shared context kept “clinician-safe” behavior aligned across frontend and backend.",
            "MCPs (Slack, Notion, etc.) pulled judging criteria and feedback in without tab sprawl.",
            "Skills and hooks for repeatable checks before we record the dashboard demo.",
        ],
    )

    add_bullets(
        prs,
        "Impact",
        [
            "Non-responders surface in the same UI used for exploration—no context switch.",
            "Chat grounded in live app_context means answers respect the cohort she just filtered.",
            "AgentCore narrative: this assistant pattern is deployable, not just a laptop prototype.",
        ],
    )

    add_bullets(
        prs,
        "Live demo & next steps",
        [
            "Follow Dr. Alvarez: Data → Explore (flag outliers) → Patient Matcher → Chat (impactful features).",
            "Insert your screenshot slide here after recording the walkthrough.",
            "Next: outcome gold standards, prospective validation, IRB-aligned governance.",
            "Thank you — First Non-responders",
        ],
    )

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
